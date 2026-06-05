import io
import base64
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from models import ClientData, Policy, AdvisorInfo, AnalysisResult

# Try Chinese fonts in order of availability (macOS, Linux)
matplotlib.rcParams["font.family"] = [
    "Arial Unicode MS", "PingFang TC", "Heiti TC",
    "Noto Sans CJK TC", "WenQuanYi Micro Hei", "sans-serif"
]
matplotlib.rcParams["axes.unicode_minus"] = False

NAVY  = RGBColor(0x0D, 0x2E, 0x5A)
TEAL  = RGBColor(0x0B, 0x7A, 0x75)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED   = RGBColor(0xDC, 0x26, 0x26)
YELLOW = RGBColor(0xD9, 0x77, 0x06)
GREEN  = RGBColor(0x05, 0x96, 0x69)
LIGHT_BLUE = RGBColor(0xF0, 0xF9, 0xFF)
BORDER_BLUE = RGBColor(0xBA, 0xE6, 0xFD)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
DARK_GRAY = RGBColor(0x64, 0x74, 0x8B)


def _add_text(slide, text, left, top, width, height,
              size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
              wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color if line_color else fill_color
    return shape


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Slide 1: Cover ────────────────────────────────────────────────────────────

def _slide1_cover(slide, client: ClientData, advisor: AdvisorInfo):
    _bg(slide, NAVY)
    today_str = date.today().strftime("%Y 年 %m 月 %d 日")

    _add_text(slide, "保障規劃分析報告",
              1, 1.1, 11.33, 1.2, size=36, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, f"客戶：{client.name}",
              1, 2.6, 11.33, 0.8, size=24,
              color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)
    _add_text(slide, f"分析日期：{today_str}",
              1, 3.3, 11.33, 0.6, size=16,
              color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)

    # Teal divider
    _add_rect(slide, 1.5, 4.1, 10.33, 0.04, TEAL)

    # Advisor card
    _add_text(slide, "您的專屬顧問",
              1, 4.3, 11.33, 0.45, size=13,
              color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)
    _add_text(slide, advisor.name,
              1, 4.7, 11.33, 0.7, size=22, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, f"{advisor.company}　{advisor.unit}",
              1, 5.35, 11.33, 0.5, size=14,
              color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)
    _add_text(slide, f"手機：{advisor.phone}　LINE：{advisor.line_id}",
              1, 5.8, 11.33, 0.5, size=12,
              color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)

    # Advisor photo
    if advisor.photo_base64:
        try:
            img_data = advisor.photo_base64
            if "," in img_data:
                img_data = img_data.split(",", 1)[1]
            img_bytes = base64.b64decode(img_data)
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(6.0), Inches(4.25), Inches(1.0), Inches(1.0)
            )
        except Exception:
            pass


# ── Slide 2: Coverage Summary ─────────────────────────────────────────────────

def _slide2_summary(slide, client: ClientData, analysis: AnalysisResult):
    _bg(slide, WHITE)
    _add_rect(slide, 0, 0, 13.33, 1.05, NAVY)
    _add_text(slide, "保障總覽", 0.3, 0.2, 8, 0.65, size=24, bold=True)
    _add_text(slide, f"客戶：{client.name}", 9.5, 0.25, 3.5, 0.55,
              size=14, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.RIGHT)

    cs = analysis.coverage_summary
    cards = [
        ("壽險保額",    f"{cs.life_total / 10000:.0f} 萬元"),
        ("醫療日額",    f"{cs.medical_daily:,.0f} 元"),
        ("癌症一次金",  f"{cs.cancer_lump_sum / 10000:.0f} 萬元"),
        ("失能月給付",  f"{cs.disability_monthly:,.0f} 元"),
        ("長照保障",    "已規劃" if cs.long_care_planned else "未規劃"),
        ("意外保額",    f"{cs.accident_total / 10000:.0f} 萬元"),
        ("年繳保費合計", f"{cs.total_annual_premium:,.0f} 元"),
    ]

    cw, ch = 3.1, 1.3
    for i, (label, val) in enumerate(cards):
        col, row = i % 4, i // 4
        x = 0.3 + col * (cw + 0.08)
        y = 1.2 + row * (ch + 0.15)
        _add_rect(slide, x, y, cw, ch, LIGHT_BLUE, BORDER_BLUE)
        _add_text(slide, label, x + 0.1, y + 0.08, cw - 0.2, 0.4,
                  size=11, color=DARK_GRAY)
        _add_text(slide, val, x + 0.1, y + 0.55, cw - 0.2, 0.6,
                  size=19, bold=True, color=NAVY)


# ── Slide 3: Policy Details ───────────────────────────────────────────────────

def _slide3_details(slide, policies: list[Policy]):
    _bg(slide, WHITE)
    _add_rect(slide, 0, 0, 13.33, 1.05, NAVY)
    _add_text(slide, "各公司保單細項明細", 0.3, 0.2, 12, 0.65, size=24, bold=True)

    headers = ["保險公司", "險種", "商品名稱", "保額", "年繳保費", "保費型態", "保障至"]
    widths  = [2.1, 1.4, 2.6, 1.55, 1.55, 1.55, 1.2]
    start_x, y = 0.15, 1.15

    x = start_x
    for h, w in zip(headers, widths):
        _add_rect(slide, x, y, w, 0.42, TEAL)
        _add_text(slide, h, x + 0.05, y + 0.06, w - 0.1, 0.32,
                  size=10, bold=True, align=PP_ALIGN.CENTER)
        x += w

    for i, p in enumerate(policies[:8]):
        y_r = 1.57 + i * 0.52
        bg = RGBColor(0xF8, 0xFA, 0xFF) if i % 2 == 0 else WHITE
        row_vals = [
            p.company, p.insurance_type, p.product_name,
            f"{p.coverage_amount / 10000:.0f}萬" if p.coverage_amount >= 10000 else f"{p.coverage_amount:,.0f}",
            f"{p.annual_premium:,.0f}",
            p.premium_type, f"{p.coverage_end_age}歲"
        ]
        x = start_x
        for val, w in zip(row_vals, widths):
            _add_rect(slide, x, y_r, w, 0.48, bg, RGBColor(0xE5, 0xE7, 0xEB))
            _add_text(slide, val, x + 0.05, y_r + 0.1, w - 0.1, 0.3,
                      size=9, color=NAVY)
            x += w


# ── Slide 4: Premium Trend ────────────────────────────────────────────────────

def _slide4_trend(slide, analysis: AnalysisResult):
    _bg(slide, WHITE)
    _add_rect(slide, 0, 0, 13.33, 1.05, NAVY)
    _add_text(slide, "自然保費趨勢分析", 0.3, 0.2, 12, 0.65, size=24, bold=True)

    trend = analysis.premium_trend
    ages    = [d.age for d in trend]
    natural = [d.natural_premium for d in trend]
    level   = [d.level_premium for d in trend]

    fig, ax = plt.subplots(figsize=(11, 4.8))
    ax.plot(ages, natural, color="#0D2E5A", marker="o", linewidth=2.5,
            label="自然保費（估算）", markersize=5)
    ax.fill_between(ages, natural, alpha=0.12, color="#0D2E5A")
    ax.plot(ages, level, color="#0B7A75", marker="s", linewidth=2,
            linestyle="--", label="平準保費（固定）", markersize=5)
    ax.set_xlabel("年齡", fontsize=12)
    ax.set_ylabel("年繳保費（元）", fontsize=12)
    ax.legend(fontsize=11, loc="upper left")
    ax.grid(True, alpha=0.3)
    ax.yaxis.set_major_formatter(
        matplotlib.ticker.FuncFormatter(lambda v, _: f"{v:,.0f}")
    )
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.7))


# ── Slide 5: Gap Analysis ─────────────────────────────────────────────────────

def _slide5_gaps(slide, analysis: AnalysisResult):
    _bg(slide, WHITE)
    _add_rect(slide, 0, 0, 13.33, 1.05, NAVY)
    _add_text(slide, "保障缺口分析", 0.3, 0.2, 12, 0.65, size=24, bold=True)

    status_map = {
        "足夠":   (GREEN,  "✓"),
        "偏低":   (YELLOW, "!"),
        "嚴重不足": (RED,    "✕"),
    }

    for i, gap in enumerate(analysis.gap_analysis[:5]):
        y = 1.2 + i * 1.2
        color, icon = status_map.get(gap.status, (YELLOW, "!"))

        _add_rect(slide, 0.25, y, 0.22, 0.85, color)
        _add_text(slide, icon, 0.25, y + 0.18, 0.22, 0.5,
                  size=14, bold=True, align=PP_ALIGN.CENTER)

        _add_text(slide, gap.category, 0.6, y, 2.0, 0.48,
                  size=16, bold=True, color=NAVY)
        _add_rect(slide, 2.65, y + 0.05, 1.3, 0.38, color)
        _add_text(slide, gap.status, 2.65, y + 0.05, 1.3, 0.38,
                  size=10, bold=True, align=PP_ALIGN.CENTER)

        _add_text(slide, f"現有：{gap.current_amount:,.0f}", 4.1, y, 3.0, 0.42,
                  size=11, color=DARK_GRAY)
        _add_text(slide, f"建議：{gap.recommended_amount:,.0f}", 7.2, y, 3.5, 0.42,
                  size=11, color=TEAL)
        _add_text(slide, gap.description, 0.6, y + 0.5, 12.4, 0.55,
                  size=10, color=GRAY_TEXT)


# ── Slide 6: Recommendations ──────────────────────────────────────────────────

def _slide6_recommendations(slide, analysis: AnalysisResult, advisor: AdvisorInfo):
    _bg(slide, WHITE)
    _add_rect(slide, 0, 0, 13.33, 1.05, NAVY)
    _add_text(slide, "加保建議與下一步", 0.3, 0.2, 12, 0.65, size=24, bold=True)

    for i, rec in enumerate(analysis.recommendations[:5]):
        y = 1.2 + i * 1.05
        badge_color = TEAL if i == 0 else NAVY
        _add_rect(slide, 0.3, y, 0.55, 0.55, badge_color)
        _add_text(slide, str(rec.priority), 0.3, y + 0.05, 0.55, 0.45,
                  size=15, bold=True, align=PP_ALIGN.CENTER)

        _add_text(slide, rec.category, 1.0, y, 2.5, 0.5,
                  size=15, bold=True, color=NAVY)
        _add_text(slide, f"建議保額：{rec.recommended_amount:,.0f} 元",
                  3.65, y, 4.0, 0.5, size=12, color=TEAL)
        _add_text(slide, rec.description, 1.0, y + 0.55, 12.0, 0.42,
                  size=10, color=GRAY_TEXT)

    # Footer
    _add_rect(slide, 0, 6.38, 13.33, 1.12, LIGHT_BLUE, BORDER_BLUE)
    _add_text(
        slide,
        f"聯絡您的顧問：{advisor.name}　|　{advisor.company} {advisor.unit}"
        f"　|　{advisor.phone}　|　LINE：{advisor.line_id}",
        0.3, 6.6, 12.73, 0.6, size=11, color=NAVY, align=PP_ALIGN.CENTER
    )


# ── Main entry point ──────────────────────────────────────────────────────────

def generate_pptx(
    client: ClientData,
    policies: list[Policy],
    advisor: AdvisorInfo,
    analysis: AnalysisResult
) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    _slide1_cover(prs.slides.add_slide(blank), client, advisor)
    _slide2_summary(prs.slides.add_slide(blank), client, analysis)
    _slide3_details(prs.slides.add_slide(blank), policies)
    _slide4_trend(prs.slides.add_slide(blank), analysis)
    _slide5_gaps(prs.slides.add_slide(blank), analysis)
    _slide6_recommendations(prs.slides.add_slide(blank), analysis, advisor)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
